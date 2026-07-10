"""PPTX parser using Docling with python-pptx fallback."""
from __future__ import annotations

import importlib.util
from typing import Any

from build_your_own_rag.models import (
    DocumentMetadata,
    ExtractionStrategy,
    ParsedDocument,
    ParsedPage,
    SourceDocument,
    utc_now_iso,
)
from build_your_own_rag.utils.logger import get_logger

logger = get_logger(__name__)


class PptxParseError(RuntimeError):
    """Raised when PPTX parsing fails."""


def parse_pptx(source: SourceDocument, strategy: ExtractionStrategy | None = None) -> ParsedDocument:
    # Use FAST as default if no strategy is provided
    strategy = strategy or ExtractionStrategy.FAST
    parser_name = "docling"
    processing_started_at = utc_now_iso()
    docling_text: str | None = None
    docling_error: str | None = None
    status = "success"

    try:
        logger.info("Attempting Docling PPTX parse", extra={"extra_info": {"source_id": source.source_id}})
        docling_text = _parse_with_docling(source)
        logger.info("Docling PPTX parse successful", extra={"extra_info": {"source_id": source.source_id}})
    except Exception as exc:
        docling_error = str(exc)
        parser_name = "python_pptx_fallback"
        status = "fallback"
        logger.warning(
            "Docling PPTX parse failed, falling back to python-pptx",
            exc_info=exc,
            extra={"extra_info": {"source_id": source.source_id, "error": docling_error}}
        )

    pptx_metadata: dict[str, Any] = {}
    pages: list[ParsedPage] = []
    
    if status == "fallback" or docling_text is None:
        try:
            logger.info("Extracting PPTX text with python-pptx", extra={"extra_info": {"source_id": source.source_id}})
            pages, pptx_metadata = _extract_with_python_pptx(source)
        except Exception as exc:
            logger.error(
                "Failed to extract PPTX text with python-pptx",
                exc_info=exc,
                extra={"extra_info": {"source_id": source.source_id, "path": str(source.path)}}
            )
            raise PptxParseError(f"Failed to extract PPTX text from {source.path}: {exc}") from exc

    text = (docling_text or "\n\n".join(page.text for page in pages)).strip()
    processing_finished_at = utc_now_iso()

    metadata = DocumentMetadata(
        core={
            "filename": source.filename,
            "extension": source.extension,
            "source_path": str(source.path),
            "size_bytes": source.size_bytes,
            "mime_type": source.mime_type,
            "source_id": source.source_id,
            "source_type": source.source_type,
        },
        processing={
            "parser_name": parser_name,
            "requested_parser": "docling",
            "strategy": strategy.value,
            "ocr_used": False,
            "processing_status": status,
            "started_at": processing_started_at,
            "finished_at": processing_finished_at,
            "docling_error": docling_error,
        },
        format_specific={
            "pptx": pptx_metadata,
        },
        elements=[
            {
                "page_number": page.page_number,
                "section_hierarchy": None,
                "coordinates": None,
                "style_name": None,
                "formatting": None,
                "language": "en",
            }
            for page in pages
        ],
    )

    return ParsedDocument(
        source=source,
        text=text,
        pages=pages,
        metadata=metadata,
        parser_name=parser_name,
        strategy=strategy,
        ocr_used=False,
        status=status,
    )


def _parse_with_docling(source: SourceDocument) -> str:
    try:
        from docling.datamodel.base_models import InputFormat
        from docling.document_converter import DocumentConverter
    except Exception as exc:
        raise RuntimeError("Docling is not installed.") from exc

    converter = DocumentConverter(allowed_formats=[InputFormat.PPTX])
    result = converter.convert(str(source.path))
    document = result.document

    if hasattr(document, "export_to_markdown"):
        return document.export_to_markdown()
    if hasattr(document, "export_to_text"):
        return document.export_to_text()
    return str(document)


def _extract_with_python_pptx(source: SourceDocument) -> tuple[list[ParsedPage], dict[str, Any]]:
    if importlib.util.find_spec("pptx") is None:
        raise RuntimeError("python-pptx is not installed. Install requirements.txt or run `python -m pip install python-pptx`.")
    
    import pptx
    prs = pptx.Presentation(str(source.path))
    
    pages = []
    
    for i, slide in enumerate(prs.slides, start=1):
        text_blocks = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_blocks.append(shape.text.strip())
            
            # Extract text from tables if the shape is a table
            if shape.has_table:
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_data:
                        text_blocks.append(" | ".join(row_data))
                        
        page_text = "\n\n".join(text_blocks)
        pages.append(
            ParsedPage(
                page_number=i,
                text=page_text,
                metadata={"page_number": i},
            )
        )
    
    pptx_metadata = {
        "slide_count": len(prs.slides),
    }
    
    return pages, pptx_metadata


def register() -> None:
    from build_your_own_rag.parsing.parser_registry import register_parser
    register_parser(".pptx", parse_pptx)
