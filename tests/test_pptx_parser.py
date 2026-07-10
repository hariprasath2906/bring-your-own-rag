import pytest
import pptx
from pathlib import Path

from build_your_own_rag.models import SourceDocument, ExtractionStrategy
from build_your_own_rag.parsing.pptx_parser import parse_pptx

@pytest.fixture
def sample_pptx_file(tmp_path):
    # Create a real pptx file using python-pptx
    file_path = tmp_path / "test.pptx"
    prs = pptx.Presentation()
    
    # Add a blank slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add a text box
    txBox = slide.shapes.add_textbox(left=0, top=0, width=100, height=100)
    tf = txBox.text_frame
    tf.text = "This is a test PPTX file."
    
    # Add a table
    table_shape = slide.shapes.add_table(rows=2, cols=2, left=0, top=100, width=100, height=100)
    table = table_shape.table
    table.cell(0, 0).text = 'Header 1'
    table.cell(0, 1).text = 'Header 2'
    table.cell(1, 0).text = 'Value 1'
    table.cell(1, 1).text = 'Value 2'
    
    prs.save(file_path)
    return file_path


def test_parse_pptx_success(sample_pptx_file):
    source = SourceDocument(
        source_id="test_id_pptx",
        source_type="local_pptx",
        path=sample_pptx_file,
        filename="test.pptx",
        extension="pptx",
        size_bytes=sample_pptx_file.stat().st_size,
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    
    parsed = parse_pptx(source, strategy=ExtractionStrategy.FAST)
    
    assert parsed.parser_name == "docling"
    assert parsed.source == source
    
    assert "This is a test PPTX file." in parsed.text
    
    # Check that it extracted the table
    assert "Header 1" in parsed.text
    assert "Value 2" in parsed.text
    
    assert "pptx" in parsed.metadata.format_specific


def test_parse_pptx_fallback(monkeypatch, sample_pptx_file):
    source = SourceDocument(
        source_id="test_id_pptx_fallback",
        source_type="local_pptx",
        path=sample_pptx_file,
        filename="test.pptx",
        extension="pptx",
        size_bytes=sample_pptx_file.stat().st_size,
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    
    # Force Docling to fail to trigger fallback
    def mock_parse_with_docling(*args, **kwargs):
        raise RuntimeError("Mock Docling failure")
    
    import build_your_own_rag.parsing.pptx_parser as pptx_module
    monkeypatch.setattr(pptx_module, "_parse_with_docling", mock_parse_with_docling)
    
    parsed = parse_pptx(source, strategy=ExtractionStrategy.FAST)
    
    assert parsed.parser_name == "python_pptx_fallback"
    assert parsed.status == "fallback"
    
    assert "This is a test PPTX file." in parsed.text
    assert "Header 1" in parsed.text
    assert "Value 2" in parsed.text
    
    pptx_meta = parsed.metadata.format_specific["pptx"]
    assert pptx_meta["slide_count"] == 1
